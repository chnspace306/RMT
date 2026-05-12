"""
生成可编辑的毕业答辩 PPTX — RMT Analytics 随机矩阵理论特征值谱分析平台
"""
import os, sys
os.system(f"{sys.executable} -m pip install python-pptx --quiet")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色方案 ──
BG_DARK   = RGBColor(0x0B, 0x0B, 0x1A)  # 深蓝黑背景
ACCENT    = RGBColor(0x00, 0xBF, 0xD8)  # 青蓝强调
ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)  # 紫色强调
ACCENT3   = RGBColor(0x10, 0xB9, 0x81)  # 绿色强调
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)
GRAY      = RGBColor(0x88, 0x88, 0x99)
DARK_CARD = RGBColor(0x15, 0x15, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

# ── 工具函数 ──
def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_decorated_title(slide, title_text, subtitle_text=None):
    """统一的标题栏装饰"""
    # 顶部装饰线
    add_box(slide, Inches(0), Inches(0), SW, Pt(4), ACCENT)
    # 标题
    add_text(slide, Inches(1), Inches(0.4), Inches(11), Inches(0.7),
             title_text, font_size=32, color=WHITE, bold=True)
    # 分隔线
    add_box(slide, Inches(1), Inches(1.15), Inches(2.5), Pt(3), ACCENT)
    if subtitle_text:
        add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.4),
                 subtitle_text, font_size=14, color=GRAY)

def add_card(slide, left, top, width, height, title, items, color=ACCENT):
    """内容卡片"""
    card = add_box(slide, left, top, width, height, DARK_CARD)
    # 左边色条
    add_box(slide, left, top, Pt(4), height, color)
    # 标题
    add_text(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
             title, font_size=16, color=color, bold=True)
    # 条目
    y = top + Inches(0.6)
    for item in items:
        add_text(slide, left + Inches(0.4), y, width - Inches(0.6), Inches(0.35),
                 f"▸ {item}", font_size=13, color=LIGHT)
        y += Inches(0.35)


# ═══════════════════════════════════════
# 第 1 页：封面
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
# 顶部大装饰条
add_box(s, Inches(0), Inches(0), SW, Inches(0.35), ACCENT)
# 主标题
add_text(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
         "基于随机矩阵理论的特征值谱分析平台", font_size=44, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
# 英文副标题
add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(0.6),
         "RMT Analytics — Eigenvalue Spectrum Analysis Platform", font_size=20, color=ACCENT,
         alignment=PP_ALIGN.CENTER)
# 分隔装饰
add_box(s, Inches(5), Inches(3.8), Inches(3.3), Pt(3), ACCENT2)
# 信息区
add_text(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
         "毕业设计答辩", font_size=24, color=LIGHT, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.1), Inches(11.3), Inches(0.4),
         "计算机科学与技术学院 · 2026年5月", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
# 底部装饰
add_box(s, Inches(0), Inches(7.15), SW, Inches(0.35), ACCENT)


# ═══════════════════════════════════════
# 第 2 页：目录
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "汇报大纲", "CONTENTS")

toc_items = [
    ("01", "研究背景与意义", "为什么选择随机矩阵理论？"),
    ("02", "RMT 核心理论", "Wigner 半圆律 · Marchenko-Pastur 分布"),
    ("03", "系统架构设计", "FastAPI + Vue 3 全栈架构"),
    ("04", "核心算法与功能", "信号检测 · RMT降噪 · IPR · 滚动窗口"),
    ("05", "四大功能视图", "Spectrum · IPR · Heatmap · Rolling Risk"),
    ("06", "技术实现与部署", "Docker · Nginx · Zeabur 云端部署"),
    ("07", "应用场景演示", "金融 · 环境 · 神经科学"),
    ("08", "总结与展望", "创新点 · 未来工作"),
]
y = Inches(1.8)
for num, title, desc in toc_items:
    # 编号圆圈
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Pt(2), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Arial"
    # 标题
    add_text(s, Inches(2.3), y, Inches(5), Inches(0.35), title, font_size=20, color=WHITE, bold=True)
    add_text(s, Inches(2.3), y + Inches(0.3), Inches(8), Inches(0.3), desc, font_size=13, color=GRAY)
    y += Inches(0.75)


# ═══════════════════════════════════════
# 第 3 页：研究背景与意义
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "研究背景与意义", "BACKGROUND & MOTIVATION")

# 左卡片 - 问题
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5),
         "核心问题", [
             "高维数据中如何区分真实信号与随机噪声？",
             "传统 PCA 需要主观设定主成分数量",
             "金融协方差矩阵包含大量统计噪声",
             "缺乏统一的异常信号检测理论框架",
         ], color=ACCENT2)

# 右卡片 - RMT 方案
add_card(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(4.5),
         "RMT 提供的解决方案", [
             "M-P 分布给出特征值的理论边界 λ₊ 和 λ₋",
             "超出 λ₊ 即标记为「统计显著」的信号",
             "无需人工设定阈值，完全由数据维度 q=p/n 决定",
             "可用于降噪、信号提取、系统性风险监测",
         ], color=ACCENT3)

# 底部关键句
add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
         "核心思想：利用随机矩阵的普适性统计规律，为高维数据分析提供一个非参数化的信号检测框架。",
         font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# 第 4 页：RMT 核心理论
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "RMT 核心理论", "WIGNER SEMICIRCLE LAW & MARCHENKO-PASTUR DISTRIBUTION")

# Wigner
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.2),
         "Wigner 半圆律", [
             "适用对象：高斯正交系 (GOE) 随机矩阵",
             "H = (A+Aᵀ)/√2，A ∼ N(0,σ²)",
             "ρ(λ) = (2/πR²)·√(R²−λ²)，|λ|≤R=2σ√N",
             "揭示了随机对称矩阵特征值分布的普适性规律",
         ], color=ACCENT)

# M-P
add_card(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(3.2),
         "Marchenko-Pastur 分布", [
             "适用对象：样本协方差矩阵 C = XᵀX/n",
             "X 为 n×p 随机矩阵，定义 q = p/n",
             "λ₊ = σ²(1+√q)²    λ₋ = σ²(1−√q)²",
             "ρ(λ) = √[(λ₊−λ)(λ−λ₋)] / (2πqσ²λ)",
         ], color=ACCENT2)

# 信号检测原理
add_card(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.8),
         "信号检测原理", [
             "任意 λ > λ₊ 的特征值在统计上极不可能由纯随机噪声产生",
             "λ₊ 提供了一个完全由数据维度决定的自动阈值——无需人工调参",
             "特征值超出 λ₊ 的数量即真实信号成分的估计数量",
         ], color=ACCENT3)


# ═══════════════════════════════════════
# 第 5 页：系统架构设计
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "系统架构设计", "SYSTEM ARCHITECTURE")

add_card(s, Inches(0.8), Inches(1.9), Inches(3.7), Inches(2.5),
         "前端层", [
             "Vue 3 + TypeScript",
             "ECharts 可视化",
             "Tailwind CSS 暗色主题",
             "marked + KaTeX AI报告渲染",
         ], color=ACCENT)

add_card(s, Inches(4.8), Inches(1.9), Inches(3.7), Inches(2.5),
         "网关层", [
             "Nginx 反向代理",
             "/api/* → 后端服务",
             "静态资源 1年缓存",
             "Gzip 压缩",
         ], color=ACCENT2)

add_card(s, Inches(8.8), Inches(1.9), Inches(3.7), Inches(2.5),
         "后端层", [
             "FastAPI + Uvicorn",
             "NumPy/SciPy 数值计算",
             "Pandas 数据处理",
             "OpenAI SDK SSE流式AI",
         ], color=ACCENT3)

# 底部流程图风格
add_box(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.5), DARK_CARD)
flow_text = (
    "数据流：CSV上传 → Pandas预处理(缺失值/标准化) → 协方差矩阵 C=XᵀX/n\n"
    "        → eigh(C)特征值分解 → MP理论边界计算 → 异常值检测\n"
    "        → IPR局部化分析 → RMT降噪热力图 → 特征向量成分提取 → AI分析"
)
txBox = s.shapes.add_textbox(Inches(1.2), Inches(4.9), Inches(10.8), Inches(2))
tf = txBox.text_frame; tf.word_wrap = True
for i, line in enumerate(flow_text.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14) if i == 0 else Pt(12)
    p.font.color.rgb = ACCENT if i == 0 else LIGHT
    p.font.name = "Consolas" if i > 0 else "Microsoft YaHei"
    p.font.bold = (i == 0)


# ═══════════════════════════════════════
# 第 6 页：核心算法流程
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "核心分析引擎 — 9步处理流水线", "CORE ANALYSIS PIPELINE")

steps = [
    ("1", "数据清洗", "移除全NaN行列，计算稀疏度"),
    ("2", "缺失值填充", "零填充 / 列均值 / 删除行"),
    ("3", "标准化(可选)", "z-score标准化或保留原始协方差"),
    ("4", "协方差矩阵", "C = (1/n)·XᵀX"),
    ("5", "特征值分解", "np.linalg.eigh(C) —— 对称矩阵优化"),
    ("6", "理论边界", "λ₊=σ²(1+√q)², λ₋=σ²(1−√q)²"),
    ("7", "异常检测", "λ>λ₊ → 提取Top-5特征向量成分"),
    ("8", "IPR+降噪", "IPR=Σv⁴, 信号子空间C_clean重建"),
    ("9", "结果输出", "MPResponse: 14个字段完整返回"),
]
x, y = Inches(0.8), Inches(1.9)
for i, (num, title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    cx = x + Inches(col * 4.1)
    cy = y + Inches(row * 1.85)
    # 卡片
    card = add_box(s, cx, cy, Inches(3.7), Inches(1.6), DARK_CARD)
    # 编号
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.15), cy + Inches(0.12), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf2 = circle.text_frame; tf2.paragraphs[0].text = num
    tf2.paragraphs[0].font.size = Pt(12); tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].font.bold = True; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, cx + Inches(0.7), cy + Inches(0.15), Inches(2.8), Inches(0.35),
             title, font_size=15, color=ACCENT, bold=True)
    add_text(s, cx + Inches(0.3), cy + Inches(0.65), Inches(3.2), Inches(0.8),
             desc, font_size=12, color=LIGHT)


# ═══════════════════════════════════════
# 第 7 页：四大功能视图
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "四大功能视图", "FOUR VIEW MODES")

views = [
    ("Spectrum 特征值谱", ACCENT, [
        "经验直方图 vs 理论分布曲线叠加",
        "自适应分箱 + 异常值自动着色",
        "图表↔面板双向联动",
        "点击异常柱形→展开特征向量详情",
    ]),
    ("IPR 局部化分析", ACCENT2, [
        "IPRⱼ = Σᵢ(vᵢⱼ)⁴ 散点图",
        "度量特征模式的局域化程度",
        "IPR≈1/N 延展态, IPR≈1 局域态",
        "物理/金融系统的相变检测",
    ]),
    ("Heatmap 降噪矩阵", ACCENT3, [
        "原始 vs RMT降噪 并排对比",
        "Matplotlib后端渲染PNG",
        "信号子空间保留,噪声均值替换",
        "对角线归一化→相关系数矩阵",
    ]),
    ("Rolling 风险演化", RGBColor(0xFF, 0x45, 0x3A), [
        "滑动窗口λ₁时序折线图",
        "实时诊断结构性变异",
        "ECharts DataZoom时间轴缩放",
        "金融系统性风险/环境突变监测",
    ]),
]
y = Inches(1.9)
for title, color, items in views:
    add_card(s, Inches(0.8), y, Inches(11.5), Inches(1.2), title, items, color)
    y += Inches(1.35)


# ═══════════════════════════════════════
# 第 8 页：前端技术实现
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "前端技术实现", "FRONTEND IMPLEMENTATION")

add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.4),
         "组件架构 (Vue 3 Composition API)", [
             "App.vue — 主控组件(1051行), 30+响应式状态",
             "RmtChart.vue — ECharts直方图, 自适应分箱算法",
             "EigenvectorChart.vue — 特征向量全景Loading图",
             "rmt.ts — 6个API函数, 三态基础路径适配",
         ], color=ACCENT)

add_card(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(2.4),
         "交互设计亮点", [
             "图表↔面板双向联动 (点击异常→展开详情)",
             "AI 分析侧边抽屉 + SSE 流式 Markdown 渲染",
             "四领域语义解释系统 (通用/金融/环境/神经科学)",
             "数据集缓存 + 多CSV切换 + 中英文双语",
         ], color=ACCENT2)

add_card(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5),
         "API 基础路径的三态适配", [
             "开发模式: VITE_API_BASE_URL 或默认 http://127.0.0.1:8000",
             "自定义后端: .env 文件配置任意 VITE_API_BASE_URL",
             "生产部署: 空字符串 '', 由 Nginx 反向代理接管 /api/ 路由",
             "同一份代码无需任何修改, 即可在三种环境下正确运行",
         ], color=ACCENT3)


# ═══════════════════════════════════════
# 第 9 页：部署架构
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "部署架构", "DEPLOYMENT & DEVOPS")

add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.2),
         "Docker 多阶段构建", [
             "前端: node:20-alpine 构建 → nginx:stable-alpine 运行",
             "后端: python:3.12-slim 单阶段, pip install + uvicorn",
             "Matplotlib Agg 后端保证无头服务器兼容",
             "多路径数据集搜索适配任意部署环境",
         ], color=ACCENT)

add_card(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(2.2),
         "Zeabur 云端配置", [
             "前后端统一 8080 端口",
             "内网服务发现: rmt.zeabur.internal:8080",
             "Nginx API反向代理 + SPA路由 try_files",
             "Gzip压缩 + 静态资源1年缓存 + immutable",
         ], color=ACCENT2)

# 架构图（文字描述）
add_box(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.7), DARK_CARD)
arch_text = (
    "                          用户浏览器\n"
    "                              │\n"
    "              ┌───────────────┴───────────────┐\n"
    "              │    Nginx :8080 (前端容器)     │\n"
    "              │  /         → dist 静态文件    │\n"
    "              │  /api/*    → 反向代理 ────────┼──────┐\n"
    "              └───────────────────────────────┘      │\n"
    "                                                     ▼\n"
    "                              ┌───────────────────────────────┐\n"
    "                              │  FastAPI :8080 (后端容器)     │\n"
    "                              │  rmt.zeabur.internal:8080     │\n"
    "                              └───────────────────────────────┘"
)
txBox = s.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.8), Inches(2.5))
tf = txBox.text_frame; tf.word_wrap = True
for i, line in enumerate(arch_text.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT if "Nginx" in line or "FastAPI" in line else LIGHT
    p.font.name = "Consolas"
    p.font.bold = ("Nginx" in line or "FastAPI" in line)


# ═══════════════════════════════════════
# 第 10 页：应用场景
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "应用场景与领域语义", "APPLICATION SCENARIOS")

apps = [
    ("金融市场", ACCENT2, [
        "λ₁ 异常 → 「市场模式」系统性风险",
        "次级异常 → 行业板块效应",
        "滚动窗口 → 实时金融危机预警",
        "RMT降噪 → 更稳健的投资组合优化",
    ]),
    ("环境监测", ACCENT3, [
        "异常特征值 → 突发区域性极端事件",
        "沙尘暴/工业泄漏 多站点同步异常",
        "IPR局部化 → 污染源空间定位",
        "时序演化 → 污染事件生命周期追踪",
    ]),
    ("神经科学", ACCENT, [
        "异常特征值 → 大规模脑网络同步激活",
        "认知任务下的主导神经元活动模式",
        "IPR → 功能网络模块化程度度量",
        "降噪矩阵 → 更精确的功能连接估计",
    ]),
]
x = Inches(0.8)
for title, color, items in apps:
    add_card(s, x, Inches(1.9), Inches(3.8), Inches(4.8), title, items, color)
    x += Inches(4.1)

# 底部通用说明
add_text(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
         "平台内置四领域语义解释器，用户可一键切换，自动获取领域专业解读",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# 第 11 页：项目创新点
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "项目创新点与贡献", "INNOVATIONS & CONTRIBUTIONS")

innovations = [
    ("非参数化信号检测", "基于MP理论边界的自动阈值，消除了传统PCA人工设定主成分数量的主观性",
     ACCENT),
    ("多维可视化分析", "特征值谱+IPR局部化+RMT降噪热力图+滚动窗口演化，四视图联动探索",
     ACCENT2),
    ("领域语义解释", "内置金融/环境/神经科学四领域知识库，自动为异常特征值提供专业解读",
     ACCENT3),
    ("AI 增强分析", "集成LLM流式分析(SSE)，自动生成Markdown格式的深度洞察报告",
     RGBColor(0xFF, 0x9F, 0x0A)),
    ("全栈工程实践", "FastAPI+Vue3+Docker+Nginx完整工程链路，支持多平台一键部署",
     RGBColor(0xFF, 0x45, 0x3A)),
]
y = Inches(1.9)
for title, desc, color in innovations:
    card = add_box(s, Inches(0.8), y, Inches(11.5), Inches(0.95), DARK_CARD)
    add_box(s, Inches(0.8), y, Pt(4), Inches(0.95), color)
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(4), Inches(0.35),
             title, font_size=18, color=color, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.5), Inches(10.8), Inches(0.35),
             desc, font_size=13, color=LIGHT)
    y += Inches(1.05)


# ═══════════════════════════════════════
# 第 12 页：总结与展望
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_decorated_title(s, "总结与展望", "SUMMARY & FUTURE WORK")

add_card(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5),
         "已完成工作", [
             "✅ 实现了 Wigner 半圆律和 M-P 分布的可视化对比",
             "✅ 构建了完整的矩阵分析引擎（9步流水线）",
             "✅ 基于MP理论实现了非参数化异常信号检测",
             "✅ 集成了IPR局部化分析和RMT降噪功能",
             "✅ 实现了滚动窗口λ₁演化风险监测",
             "✅ 完成了全栈工程化与云端部署",
         ], color=ACCENT3)

add_card(s, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.5),
         "未来工作", [
             "🔮 支持更多RMT系综 (GUE, GSE, Wishart)",
             "🔮 引入Kolmogorov-Smirnov拟合优度检验",
             "🔮 扩展到张量数据 (多维随机矩阵)",
             "🔮 增加更多领域的内置示例数据集",
             "🔮 支持用户自定义AI提示词模板",
         ], color=ACCENT2)


# ═══════════════════════════════════════
# 第 13 页：致谢
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_box(s, Inches(0), Inches(0), SW, Inches(0.35), ACCENT)

add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0),
         "感谢各位老师！", font_size=48, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
         "敬请批评指正", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_box(s, Inches(5), Inches(4.0), Inches(3.3), Pt(3), ACCENT2)

add_text(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
         "RMT Analytics — 随机矩阵理论特征值谱分析平台", font_size=18, color=LIGHT,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.4),
         "FastAPI + Vue 3 + ECharts + Docker", font_size=16, color=GRAY,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.4),
         "GitHub: github.com/chnspace306/RMT", font_size=14, color=ACCENT,
         alignment=PP_ALIGN.CENTER)

add_box(s, Inches(0), Inches(7.15), SW, Inches(0.35), ACCENT)


# ── 保存 ──
OUTPUT = r"D:\毕业论文\毕业设计\RMT_Analytics_答辩演示_可编辑.pptx"
prs.save(OUTPUT)
print(f"✅ 已生成可编辑 PPTX: {OUTPUT}")
print(f"共 {len(prs.slides)} 张幻灯片")
