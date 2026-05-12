"""
将 PDF 幻灯片的每一页提取为图片，并创建 PPTX 演示文稿。
使用方法：在终端运行 python pdf_to_pptx.py
"""

import os
import sys

# 安装依赖
os.system(f"{sys.executable} -m pip install PyMuPDF python-pptx --quiet")

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# 路径配置
PDF_PATH = r"C:\Users\X\AppData\Local\Claude-3p\local-agent-mode-sessions\61a0cb38-cf3b-4d3d-a794-a1d69c1ade4e\00000000-0000-4000-8000-000000000001\local_0412bd85-f8b5-4f06-8159-0ee85c6fcc5b\uploads\RMT_Analytics_Perfect_Final.html"
OUTPUT_PPTX = r"D:\毕业论文\毕业设计\RMT_Analytics_答辩演示.pptx"

print(f"正在打开 PDF: {PDF_PATH}")
doc = fitz.open(PDF_PATH)
total_pages = len(doc)
print(f"PDF 共 {total_pages} 页，开始转换...")

# 创建 PPTX (16:9 宽屏)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 使用空白版式 (index 6 通常是 blank)
blank_layout = prs.slide_layouts[6]

for page_num in range(total_pages):
    print(f"  处理第 {page_num + 1}/{total_pages} 页...")

    page = doc[page_num]

    # 渲染页面为高分辨率 PNG
    mat = fitz.Matrix(2.0, 2.0)  # 2x 缩放保证清晰度
    pix = page.get_pixmap(matrix=mat)

    # 保存为 PNG 到内存
    img_bytes = pix.tobytes("png")

    # 创建幻灯片并添加图片
    slide = prs.slides.add_slide(blank_layout)

    # 将图片流插入幻灯片，填满整页
    img_stream = BytesIO(img_bytes)
    slide.shapes.add_picture(
        img_stream,
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=prs.slide_height
    )

    print(f"  ✓ 第 {page_num + 1} 页完成")

doc.close()

# 保存 PPTX
prs.save(OUTPUT_PPTX)
print(f"\n✅ 完成！PPTX 已保存至: {OUTPUT_PPTX}")
print(f"共 {total_pages} 张幻灯片")
