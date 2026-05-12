"""
将 PDF 答辩幻灯片转为 PPTX
用法: python convert_pdf_to_pptx.py
"""
import os, sys

# 安装依赖
os.system(f"{sys.executable} -m pip install PyMuPDF python-pptx --quiet")

import fitz
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# 路径
PDF_FILE = r"C:\Users\X\AppData\Local\Claude-3p\local-agent-mode-sessions\61a0cb38-cf3b-4d3d-a794-a1d69c1ade4e\00000000-0000-4000-8000-000000000001\local_0412bd85-f8b5-4f06-8159-0ee85c6fcc5b\uploads\202212310125_刘俊佑_基于随机矩阵的可视化分析系统_答辩ppt.pdf"
OUTPUT  = r"D:\毕业论文\毕业设计\答辩演示.pptx"

print("打开 PDF...")
doc = fitz.open(PDF_FILE)
total = len(doc)
print(f"共 {total} 页")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for i in range(total):
    page = doc[i]
    # 2x 渲染保证清晰度
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
    img_bytes = pix.tobytes("png")

    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        BytesIO(img_bytes),
        left=Inches(0), top=Inches(0),
        width=prs.slide_width, height=prs.slide_height
    )
    print(f"  第 {i+1}/{total} 页 ✓")

doc.close()
prs.save(OUTPUT)
print(f"\n✅ 已保存: {OUTPUT}")
print(f"共 {total} 张幻灯片")
